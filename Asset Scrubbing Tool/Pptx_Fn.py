#!/usr/bin/env python
# coding: utf-8

# In[ ]:


import os
from pptx import Presentation
import pptx
import re
import FileName_Remove
import FileName_Replace
from pptx.enum.shapes import MSO_SHAPE_TYPE
import Pptx_ReviewComment_Removal


# In[ ]:


def Scrubbing_Process(Location,Keywords,Flag,pptxFileList):
    global Scrubbing_Location
    global Scrubbing_Keywords
    global Scrubbing_pptxFileList
    
    Scrubbing_Location=Location
    Scrubbing_Keywords=Keywords
    Scrubbing_pptxFileList=pptxFileList
    Flag=Flag.lower()
    
    
    if (Flag == 'replace'):
        Scrubbing_Replace()
    elif (Flag == 'remove'):
        Scrubbing_Remove()
        


# In[ ]:


def Scrubbing_Replace():
    
    Replace_List=[]
    print("Replace-")
    Replace_Keywords=Scrubbing_Keywords.split(",")
    for i in Replace_Keywords:
        Single_Keyword=i.split(":")
        Replace_List.append(Single_Keyword)
    print(Replace_List)
        
    def pptElementsReplace():
        global ppt
        global pptxListremove
        pptxListremove = []
        global thisText
        text_runs = []
        pptxListremove1 = []
        global thisText1
        
        for File in Scrubbing_pptxFileList:
            try:
                ppt = Presentation(File)
                for slide in ppt.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for key,value in Replace_List:
                                text_frame = shape.text_frame
                                text_frame.paragraphs[0].add_run()
                                cur_text = text_frame.paragraphs[0].runs[0].text
#                                 print(cur_text)
                                thisText = re.sub(key + "(?=$|[^a-z])", value, cur_text, flags=re.IGNORECASE)
#                                 print(thisText)
                                text_frame.paragraphs[0].runs[0].text = thisText


                        if not shape.has_table:
                            continue
                        tbl = shape.table
                        row_count = len(tbl.rows)
                        col_count = len(tbl.columns)
                        for r in range(0, row_count):
                            for c in range(0, col_count):
                                cell = tbl.cell(r, c)
                                paragraphs = cell.text_frame.paragraphs
                                for paragraph in paragraphs:
                                    for run in paragraph.runs:
                                        text_runs.append(run.text)
                                        for key,value in Replace_List:
                                            thisText1 = re.sub(key + "(?=$|[^a-z])", value, cur_text, flags=re.IGNORECASE)
                                            run.text = thisText1


                for slide in ppt.slides:
                    # ---only operate on group shapes---
                    group_shapes = [
                        shp for shp in slide.shapes
                        if shp.shape_type == MSO_SHAPE_TYPE.GROUP
                    ]
                    for group_shape in group_shapes:
                        for shape in group_shape.shapes:
                            if shape.has_text_frame:
                                for key,value in Replace_List:
                                    text_frame = shape.text_frame
                                    text_frame.paragraphs[0].add_run()
                                    cur_text = text_frame.paragraphs[0].runs[0].text
                                    thisText = re.sub(key + "(?=$|[^a-z])", value, cur_text, flags=re.IGNORECASE)
                                    text_frame.paragraphs[0].runs[0].text = thisText

                for slide in ppt.slides:
                        if slide.has_notes_slide:
                            for key,value in Replace_List:
                                notes_slide = slide.notes_slide
                                text_frame = notes_slide.notes_text_frame
                                thisText = re.sub(key + "(?=$|[^a-z])", value, text_frame.text, flags=re.IGNORECASE)
                                text_frame.text = thisText

                ppt.save(File)
            except Exception as e:
                print(e)
                print("Couldn't open the file :",File)

    pptElementsReplace()
    FileName_Replace.FileName_Replace(Scrubbing_Location,Scrubbing_Keywords,Scrubbing_pptxFileList)

    

def Scrubbing_Remove():
    print("Remove Keywords-")
    Remove_Keywords=Scrubbing_Keywords.split(",")
    print(Remove_Keywords)
    
        
    ###function for removing Element from pptx
    def pptElementsRemoval():
        global ppt
        global pptxListremove
        pptxListremove = []
        global thisText
        text_runs = []
        pptxListremove1 = []
        global thisText1

        for File in Scrubbing_pptxFileList:
            try:
                ppt = Presentation(File)
                for slide in ppt.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for l in Remove_Keywords:
                                text_frame = shape.text_frame
                                text_frame.paragraphs[0].add_run()
                                cur_text = text_frame.paragraphs[0].runs[0].text
                                thisText = re.sub(l + "(?=$|[^a-z])", "", cur_text, flags=re.IGNORECASE)
                                text_frame.paragraphs[0].runs[0].text = thisText



                        if not shape.has_table:
                            continue
                        tbl = shape.table
                        row_count = len(tbl.rows)
                        col_count = len(tbl.columns)
                        for r in range(0, row_count):
                            for c in range(0, col_count):
                                cell = tbl.cell(r, c)
                                paragraphs = cell.text_frame.paragraphs
                                for paragraph in paragraphs:
                                    for run in paragraph.runs:
                                        text_runs.append(run.text)
                                        for z in Remove_Keywords:
                                            thisText1 = re.sub(z + "(?=$|[^a-z])", "", run.text, flags=re.IGNORECASE)
                                            run.text = thisText1

                for slide in ppt.slides:
                    # ---only operate on group shapes---
                    group_shapes = [
                        shp for shp in slide.shapes
                        if shp.shape_type == MSO_SHAPE_TYPE.GROUP
                    ]
                    for group_shape in group_shapes:
                        for shape in group_shape.shapes:
                            if shape.has_text_frame:
                                for x in Remove_Keywords:
                                    text_frame = shape.text_frame
                                    text_frame.paragraphs[0].add_run()
                                    cur_text = text_frame.paragraphs[0].runs[0].text
                                    thisText = re.sub(x + "(?=$|[^a-z])", "", cur_text, flags=re.IGNORECASE)
                                    text_frame.paragraphs[0].runs[0].text = thisText

                for slide in ppt.slides:
                        if slide.has_notes_slide:
                            for y in Remove_Keywords:
                                notes_slide = slide.notes_slide
                                text_frame = notes_slide.notes_text_frame
                                if(text_frame!=None):
                                    thisText = re.sub(y + "(?=$|[^a-z])", "", text_frame.text, flags=re.IGNORECASE)
                                    text_frame.text = thisText


                ppt.save(os.path.abspath(File))
            except:
                print("Couldn't open the file :",File)
    pptElementsRemoval()

    FileName_Remove.FileName_Removal(Scrubbing_Location,Scrubbing_Keywords,Scrubbing_pptxFileList)
    #Pptx_ReviewComment_Removal.reviewcommentsRemovalFromPPT(Scrubbing_Location,Remove_Keywords,Scrubbing_pptxFileList)


# In[ ]:




